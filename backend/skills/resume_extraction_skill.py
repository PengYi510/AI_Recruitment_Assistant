"""简历提取Skill - 从PDF/Word/MD/TXT中提取结构化数据，写入SQL+向量库

支持格式: PDF, DOCX, MD, TXT, PPTX
提取流程:
1. 解析文件获取原始文本
2. 使用LLM进行结构化信息抽取（输出标准JSON schema）
3. 验证提取结果的完整性
4. 写入SQL 6张表 + ChromaDB向量库
"""
import logging
import json
import re
from typing import Dict, Any, Optional
from pathlib import Path
from backend.skills.base_skill import BaseSkill

logger = logging.getLogger(__name__)

# 提取结果的JSON Schema定义
RESUME_SCHEMA = {
    "personal_info": {
        "name": "str",
        "gender": "str|null",
        "birth_date": "str|null",
        "age": "int|null",
        "phone": "str|null",
        "email": "str|null",
        "address": "str|null",
        "current_city": "str|null (省份/城市，如 北京/北京、广东/深圳；海外则 国家/城市)",
        "hometown": "str|null (省份/城市，如 湖南/长沙；海外则 国家/城市)",
        "current_position": "str|null",
        "work_years": "int",
        "expected_salary": "float|null",
        "job_status": "str|null",
        "location": "str|null",
        "highest_education": "str",
        "summary": "str|null"
    },
    "education_history": [{
        "degree": "str (专科/本科/硕士/博士)",
        "school": "str",
        "major": "str",
        "start_date": "str (YYYY.MM)",
        "end_date": "str (YYYY.MM 或 至今)",
        "is_fulltime": "bool",
        "school_tier": "str|null (985/211/双一流/普通本科)"
    }],
    "skills": [{
        "skill_name": "str",
        "proficiency": "int (1-5, 1=了解 2=一般 3=熟悉 4=熟练 5=精通)",
        "category": "str|null (数据库/大数据/编程语言/AI/框架/工具)"
    }],
    "awards_certificates": [{
        "type": "str (award/certificate)",
        "name": "str",
        "level": "str|null (国家级/省级/校级/企业级)",
        "date": "str|null (YYYY.MM)",
        "role": "str|null (队长/队员/个人)",
        "description": "str|null"
    }],
    "work_experiences": [{
        "company_name": "str",
        "position": "str",
        "location": "str|null",
        "start_date": "str (YYYY.MM)",
        "end_date": "str (YYYY.MM 或 至今)",
        "duration_months": "int|null",
        "description": "str"
    }],
    "projects": [{
        "project_name": "str",
        "role": "str",
        "start_date": "str|null (YYYY.MM)",
        "end_date": "str|null (YYYY.MM)",
        "duration_months": "int|null",
        "description": "str",
        "technologies": "str (逗号分隔的技术栈)"
    }],
    "publications": [{
        "title": "str (论文标题)",
        "venue": "str (期刊/会议名称，如 IEEE TPAMI, NeurIPS)",
        "venue_type": "str (journal/conference)",
        "venue_rank": "str (CCF-A/CCF-B/CCF-C)",
        "sci_zone": "str|null (Q1/Q2/Q3/Q4，仅期刊)",
        "year": "int (发表年份)",
        "author_position": "str (first_author/corresponding_author/co_author)",
        "doi": "str|null"
    }],
    "conferences": [{
        "conference_name": "str (会议全称或缩写)",
        "conference_rank": "str (CCF-A/CCF-B/CCF-C)",
        "year": "int (参会年份)",
        "location": "str|null (举办地点)",
        "role": "str (oral_presentation/poster/workshop/attendee/session_chair/keynote)",
        "paper_title": "str|null (如有投稿论文)",
        "description": "str|null"
    }],
    "extra_attributes": {
        "gpa": "str|null (如 3.8/4.0 或 88/100)",
        "hobbies": "str|null (逗号分隔，如 篮球,游泳,阅读)",
        "target_job": "str|null (目标岗位/期望岗位)",
        "ethnicity": "str|null (民族)",
        "height_cm": "float|null (身高cm)",
        "weight_kg": "float|null (体重kg)",
        "marital_status": "str|null (婚姻状况: 未婚/已婚/离异)",
        "political_status": "str|null (政治面貌: 党员/团员/群众)",
        "languages": "str|null (语言能力，如 英语CET-6,日语N1)",
        "certifications": "str|null (职业资格证书，如 PMP,CPA)",
        "self_evaluation": "str|null (自我评价/个人优势)",
        "other": "str|null (其他未归类的信息)"
    }
}

# LLM提取Prompt模板
EXTRACTION_PROMPT = """你是一个专业的简历信息提取助手。请从以下简历文本中提取结构化信息，严格按照JSON格式输出。

## 提取规则:
1. **个人信息**: 提取姓名、性别、出生日期、年龄、电话、邮箱、地址、现居地(current_city)、家乡(hometown)、当前职位、工作年限、期望薪资、求职状态、所在城市、最高学历
   - current_city: 当前居住城市，格式为"省份/城市"（如"北京/北京"、"广东/深圳"），海外则"国家/城市"（如"美国/旧金山"）
   - hometown: 家乡/籍贯，格式同上（如"湖南/长沙"、"四川/成都"）
   - 如果简历中提到"现居"、"坐标"、"所在城市"等，提取为 current_city
   - 如果简历中提到"籍贯"、"家乡"、"户籍"等，提取为 hometown
2. **教育经历**: 按时间倒序，提取每段教育的学校、学历、专业、起止时间。学校层次根据知名度判断(985/211/双一流/普通本科/专科)；海外院校的 school_tier 填写"海外QS100"或"海外QS200"或"海外名校"
3. **技术栈**: 提取所有技术技能，根据简历中的描述判断熟练度(精通=5, 熟练/熟悉=4, 掌握=3, 了解/会使用=2, 接触=1)
4. **获奖/证书**: 提取竞赛获奖、资格证书等
5. **工作经历**: 提取每段工作的公司、职位、地点、起止时间、职责描述
6. **项目经历**: 提取每个项目的名称、角色、时间、描述、使用的技术栈
7. **论文发表**: 提取所有发表的论文，包括标题、发表期刊/会议名称、期刊类型(journal/conference)、CCF等级(CCF-A/CCF-B/CCF-C)、SCI分区(Q1-Q4，仅期刊)、发表年份、作者位次(first_author/corresponding_author/co_author)、DOI
8. **国际会议**: 提取参加的学术会议，包括会议名称、CCF等级、年份、地点、参与角色(oral_presentation/poster/workshop/attendee/session_chair/keynote)、投稿论文标题(如有)

## 注意事项:
- 日期格式统一为 YYYY.MM（如 2024.05）
- 如果是“至今”，end_date填“至今”
- work_years 根据最早工作开始时间计算
- 如果信息缺失，对应字段填null
- technologies字段用逗号分隔
- venue_rank 必须为 CCF-A/CCF-B/CCF-C 之一，根据期刊/会议的知名度判断
- 常见CCF-A期刊/会议: IEEE TPAMI, IJCV, ACM Computing Surveys, NeurIPS, ICML, CVPR, ACL, SIGMOD, OSDI等
- 常见CCF-B: Pattern Recognition, Neurocomputing, AAAI, IJCAI, ECCV, COLING, ICDE等
- author_position: 第一作者=first_author, 通讯作者=corresponding_author, 其他=co_author
9. **动态扩展属性**: 提取简历中出现的 GPA、爱好、目标岗位、民族、身高体重、婚姻状况、政治面貌、语言能力、自我评价等信息，放入 extra_attributes 对象中

## 简历原文:
```
{resume_text}
```

## 请输出标准JSON（不要输出其他内容）:
"""


class ResumeExtractionSkill(BaseSkill):
    """简历提取Skill: 解析简历文件 → LLM结构化抽取 → 写入数据库"""

    def __init__(self):
        super().__init__(name="resume_extraction", description="从简历文件中提取结构化数据")

    async def execute(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """执行简历提取

        params:
            resume_text: 直接传入的简历纯文本（优先级高于 file_path）
            file_path: 简历文件路径（resume_text 未提供时使用）
            save_to_db: 是否写入数据库（默认True）
        """
        resume_text = params.get("resume_text", "")
        file_path = params.get("file_path", "")
        save_to_db = params.get("save_to_db", True)

        # Step 1: 获取原始文本——优先使用直接传入的文本，否则从文件解析
        if resume_text and resume_text.strip():
            raw_text = resume_text
        elif file_path:
            raw_text = self._extract_text(file_path)
            if not raw_text:
                return {"success": False, "error": f"Failed to extract text from {file_path}"}
        else:
            return {"success": False, "error": "either resume_text or file_path is required"}

        # Step 2: LLM结构化抽取
        structured_data = await self._llm_extract(raw_text)
        if not structured_data:
            return {"success": False, "error": "LLM extraction failed"}

        # Step 3: 后处理——修正技能 proficiency（LLM 经常全部给默认值3）
        structured_data = self._fix_skill_proficiency(structured_data, raw_text)

        # Step 4: 验证提取结果
        validation = self._validate_extraction(structured_data)
        if not validation["valid"]:
            logger.warning(f"Extraction validation warnings: {validation['warnings']}")

        # Step 5: 写入数据库（含原始文本和动态属性）
        candidate_id = None
        if save_to_db:
            candidate_id = self._save_to_database(structured_data, raw_text)

        return {
            "success": True,
            "candidate_id": candidate_id,
            "extracted_data": structured_data,
            "validation": validation,
            "raw_text_length": len(raw_text),
            "raw_text": raw_text
        }

    def _extract_text(self, file_path: str) -> Optional[str]:
        """从各种格式的文件中提取纯文本"""
        path = Path(file_path)
        if not path.exists():
            logger.error(f"File not found: {file_path}")
            return None

        suffix = path.suffix.lower()
        try:
            if suffix == ".pdf":
                return self._extract_from_pdf(file_path)
            elif suffix == ".docx":
                return self._extract_from_docx(file_path)
            elif suffix in (".md", ".txt"):
                return path.read_text(encoding="utf-8")
            elif suffix == ".pptx":
                return self._extract_from_pptx(file_path)
            else:
                logger.error(f"Unsupported file format: {suffix}")
                return None
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return None

    def _extract_from_pdf(self, file_path: str) -> Optional[str]:
        """从PDF提取文本"""
        try:
            import pdfplumber
            texts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        texts.append(text)
            return "\n".join(texts)
        except ImportError:
            logger.error("pdfplumber not installed. Run: pip install pdfplumber")
            return None

    def _extract_from_docx(self, file_path: str) -> Optional[str]:
        """从Word文档提取文本"""
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs if para.text.strip())
        except ImportError:
            logger.error("python-docx not installed. Run: pip install python-docx")
            return None

    def _extract_from_pptx(self, file_path: str) -> Optional[str]:
        """从PPT提取文本"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text.strip())
            return "\n".join(texts)
        except ImportError:
            logger.error("python-pptx not installed. Run: pip install python-pptx")
            return None

    def _fix_skill_proficiency(self, data: Dict[str, Any], raw_text: str) -> Dict[str, Any]:
        """基于简历原文中的等级词修正技能 proficiency

        LLM 提取时经常把所有技能的 proficiency 都设为 3（默认值），
        这里通过正则匹配原文中的"精通/熟练/熟悉/掌握/了解"等关键词来修正。
        """
        skills = data.get("skills")
        if not skills or not isinstance(skills, list):
            return data

        # 检查是否所有 proficiency 都相同（说明 LLM 没有正确区分）
        proficiencies = [s.get("proficiency", 3) for s in skills if isinstance(s, dict)]
        if not proficiencies:
            return data
        all_same = len(set(proficiencies)) <= 1

        # 等级词 → proficiency 映射（从高到低匹配）
        # 注意：pattern 末尾用 [^\S\n]* 代替 \s*，避免跨行匹配
        level_patterns = [
            (5, [r"精通[^\S\n]*", r"深入掌握[^\S\n]*", r"专家级[^\S\n]*"]),
            (4, [r"熟练(?:使用|掌握|运用)?[^\S\n]*", r"熟悉[^\S\n]*", r"深入了解[^\S\n]*"]),
            (3, [r"掌握[^\S\n]*", r"能够使用[^\S\n]*", r"具备[^\n]*(?:能力|经验)[^\S\n]*"]),
            (2, [r"了解[^\S\n]*", r"会使用[^\S\n]*", r"接触过[^\S\n]*", r"基础[^\S\n]*"]),
            (1, [r"初步了解[^\S\n]*", r"入门[^\S\n]*"]),
        ]

        for skill in skills:
            if not isinstance(skill, dict):
                continue
            skill_name = skill.get("skill_name", "")
            if not skill_name:
                continue

            # 只在 proficiency 全部相同时才强制修正，否则仅修正值为3的
            if not all_same and skill.get("proficiency", 3) != 3:
                continue

            # 在原文中搜索该技能附近的等级词
            escaped_name = re.escape(skill_name)
            matched_level = None

            # 分隔符：排除中英文逗号、句号、分号、括号、换行，防止跨边界匹配
            sep = r"[^，。；,;()\（\）\n]{0,20}"

            for level, patterns in level_patterns:
                for pat in patterns:
                    # 模式1: "精通 Python"（等级词在前）
                    regex_before = pat + sep + escaped_name
                    # 模式2: "Python（精通）" 或 "Python 精通"（等级词在后）
                    regex_after = escaped_name + r"[（(\s]?" + sep + pat
                    if re.search(regex_before, raw_text) or re.search(regex_after, raw_text):
                        matched_level = level
                        break
                if matched_level:
                    break

            if matched_level:
                skill["proficiency"] = matched_level
            elif all_same:
                # 如果全部相同且没有匹配到等级词，给一个基于位置的默认分布
                # 简历中先列出的技能通常更重要/更熟练
                idx = skills.index(skill)
                total = len(skills)
                if idx < total * 0.2:
                    skill["proficiency"] = 5
                elif idx < total * 0.4:
                    skill["proficiency"] = 4
                elif idx < total * 0.7:
                    skill["proficiency"] = 3
                else:
                    skill["proficiency"] = 2

        data["skills"] = skills
        return data

    async def _llm_extract(self, raw_text: str) -> Optional[Dict[str, Any]]:
        """使用LLM进行结构化信息抽取"""
        from backend.config import LLM_API_KEY, LLM_BASE_URL, LLM_MODEL
        import httpx

        prompt = EXTRACTION_PROMPT.format(resume_text=raw_text)

        try:
            async with httpx.AsyncClient(timeout=60) as client:
                response = await client.post(
                    f"{LLM_BASE_URL}/chat/completions",
                    headers={
                        "Authorization": f"Bearer {LLM_API_KEY}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model": LLM_MODEL,
                        "messages": [{"role": "user", "content": prompt}],
                        "temperature": 0.1,
                        "max_tokens": 8192
                    }
                )
                response.raise_for_status()
                result = response.json()
                content = result["choices"][0]["message"]["content"]

                # 提取JSON（处理可能的markdown代码块包裹）
                json_str = self._extract_json_from_response(content)
                raw = json.loads(json_str)
                # 归一化 key（LLM 常输出别名，统一到标准 schema）
                return self._normalize_keys(raw)

        except Exception as e:
            logger.error(f"LLM extraction failed: {e}")
            return None

    def _normalize_keys(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """将 LLM 输出中常见的 key 别名归一化到标准 schema 字段名。

        LLM 不总是严格遵守 schema，例如会输出 education/technologies/work_experience 等，
        这里统一映射，保证后续入库逻辑能正确读取。
        """
        # 顶层 section 别名 -> 标准名
        section_alias = {
            "education": "education_history",
            "educations": "education_history",
            "education_history": "education_history",
            "edu": "education_history",
            "skills": "skills",
            "skill": "skills",
            "technologies": "skills",
            "technology": "skills",
            "tech_stack": "skills",
            "tech_skills": "skills",
            "awards_certificates": "awards_certificates",
            "awards_certifications": "awards_certificates",
            "awards": "awards_certificates",
            "certificates": "awards_certificates",
            "awards_and_certificates": "awards_certificates",
            "work_experiences": "work_experiences",
            "work_experience": "work_experiences",
            "work": "work_experiences",
            "experiences": "work_experiences",
            "projects": "projects",
            "project": "projects",
            "project_experience": "projects",
            "publications": "publications",
            "publication": "publications",
            "papers": "publications",
            "paper": "publications",
            "published_papers": "publications",
            "论文": "publications",
            "论文发表": "publications",
            "conferences": "conferences",
            "conference": "conferences",
            "conference_participation": "conferences",
            "academic_conferences": "conferences",
            "国际会议": "conferences",
            "会议": "conferences",
            "personal_info": "personal_info",
            "personal": "personal_info",
            "basic_info": "personal_info",
            "extra_attributes": "extra_attributes",
            "extra": "extra_attributes",
            "additional_info": "extra_attributes",
            "other_info": "extra_attributes",
        }
        normalized: Dict[str, Any] = {}
        for key, value in (data or {}).items():
            std = section_alias.get(key.lower().strip(), key)
            # 同一标准 key 出现多次时合并列表
            if std in normalized and isinstance(normalized[std], list) and isinstance(value, list):
                normalized[std].extend(value)
            else:
                normalized[std] = value

        # 字段级别名归一化
        self._normalize_item_fields(normalized)
        return normalized

    def _parse_skills_string(self, skills_str: str) -> list:
        """解析 LLM 返回的字符串格式技能列表

        支持格式：
        - "C++:5, Scala:4, Python:3"（技能名:proficiency）
        - "C++（精通）, Scala（熟练）, Python（掌握）"（技能名+中文等级）
        - "C++, Scala, Python"（纯技能名列表）
        - "精通C++，熟练Scala，掌握Python"（等级词+技能名）
        """
        prof_map = {"精通": 5, "熟练": 4, "熟悉": 4, "掌握": 3, "了解": 2, "接触": 1}
        results = []

        # 尝试按逗号/分号/顿号分割
        parts = re.split(r"[,，;；、\n]+", skills_str.strip())

        for part in parts:
            part = part.strip()
            if not part or len(part) <= 1:
                continue

            # 格式1: "技能名:数字" 或 "技能名：数字"
            m = re.match(r"(.+?)\s*[:：]\s*(\d)", part)
            if m:
                name = m.group(1).strip()
                prof = int(m.group(2))
                if name and len(name) > 1:
                    results.append({"skill_name": name, "proficiency": min(max(prof, 1), 5)})
                continue

            # 格式2: "技能名（等级词）" 或 "技能名(等级词)"
            m = re.match(r"(.+?)\s*[（(]\s*(精通|熟练|熟悉|掌握|了解|接触)\s*[）)]", part)
            if m:
                name = m.group(1).strip()
                prof = prof_map.get(m.group(2), 3)
                if name and len(name) > 1:
                    results.append({"skill_name": name, "proficiency": prof})
                continue

            # 格式3: "等级词+技能名"
            m = re.match(r"(精通|熟练|熟悉|掌握|了解|接触)\s*(.+)", part)
            if m:
                prof = prof_map.get(m.group(1), 3)
                name = m.group(2).strip()
                if name and len(name) > 1:
                    results.append({"skill_name": name, "proficiency": prof})
                continue

            # 格式4: 纯技能名
            if len(part) > 1:
                results.append({"skill_name": part, "proficiency": 3})

        return results if results else []

    def _normalize_item_fields(self, data: Dict[str, Any]) -> None:
        """归一化各 section 内部元素的字段别名"""
        # 技能字段：name/tech -> skill_name；level -> proficiency
        prof_map = {"精通": 5, "熟练": 4, "熟悉": 3, "掌握": 3, "了解": 2, "接触": 1}
        skills = data.get("skills")

        # 处理 LLM 返回 skills 为字符串的情况（如 "C++:5,Scala:4,Python:3"）
        if isinstance(skills, str):
            skills = self._parse_skills_string(skills)
            data["skills"] = skills

        if isinstance(skills, list):
            fixed = []
            for s in skills:
                if isinstance(s, str):
                    # 单个技能名字符串（非单字符），如 "Python"
                    if len(s) > 1:
                        fixed.append({"skill_name": s, "proficiency": 3})
                    # 单字符说明数据有问题，跳过
                elif isinstance(s, dict):
                    name = s.get("skill_name") or s.get("name") or s.get("tech") or s.get("skill")
                    prof = s.get("proficiency")
                    if prof is None:
                        lvl = s.get("level") or s.get("proficiency_level")
                        prof = prof_map.get(str(lvl).strip(), 3) if lvl else 3
                    if isinstance(prof, str):
                        prof = prof_map.get(prof.strip(), 3)
                    fixed.append({"skill_name": name, "proficiency": prof or 3,
                                  "category": s.get("category")})
            data["skills"] = [f for f in fixed if f.get("skill_name")]

        # 教育字段：degree/school/major 等已对齐，school_name -> school
        edu = data.get("education_history")
        if isinstance(edu, list):
            for e in edu:
                if isinstance(e, dict):
                    if not e.get("school") and e.get("school_name"):
                        e["school"] = e["school_name"]
                    if not e.get("degree") and e.get("education_level"):
                        e["degree"] = e["education_level"]

        # 工作经历：company/company_name
        work = data.get("work_experiences")
        if isinstance(work, list):
            for w in work:
                if isinstance(w, dict) and not w.get("company_name"):
                    w["company_name"] = w.get("company") or w.get("employer")

        # 获奖：name 已对齐；title -> name
        awards = data.get("awards_certificates")
        if isinstance(awards, list):
            for a in awards:
                if isinstance(a, dict) and not a.get("name"):
                    a["name"] = a.get("title") or a.get("award_name") or a.get("certificate_name")

        # 项目：name -> project_name
        projs = data.get("projects")
        if isinstance(projs, list):
            for p in projs:
                if isinstance(p, dict) and not p.get("project_name"):
                    p["project_name"] = p.get("name") or p.get("title")

    def _extract_json_from_response(self, content: str) -> str:
        """从LLM响应中提取JSON字符串（处理代码块包裹、JSON 后跟多余内容等情况）。

        通过括号配平截取第一个完整的 JSON 对象，避免 LLM 在 JSON 之后追加解释/
        第二个对象导致 json.loads 报 "Extra data"。
        """
        # 尝试提取```json ... ```中的内容
        json_match = re.search(r'```(?:json)?\s*\n?(.*?)\n?```', content, re.DOTALL)
        if json_match:
            return self._slice_first_json_object(json_match.group(1).strip())
        content = content.strip()
        # 从第一个 { 起，用括号配平截取首个完整对象（容忍字符串内的花括号与转义）
        sliced = self._slice_first_json_object(content)
        if sliced:
            return sliced
        # 兜底：第一个{到最后一个}
        start = content.find("{")
        end = content.rfind("}")
        if start != -1 and end != -1:
            return content[start:end + 1]
        return content

    @staticmethod
    def _slice_first_json_object(text: str) -> str:
        """从文本中按括号配平截取第一个完整的 JSON 对象，找不到返回空串。"""
        start = text.find("{")
        if start == -1:
            return ""
        depth = 0
        in_str = False
        escape = False
        for i in range(start, len(text)):
            ch = text[i]
            if in_str:
                if escape:
                    escape = False
                elif ch == "\\":
                    escape = True
                elif ch == '"':
                    in_str = False
                continue
            if ch == '"':
                in_str = True
            elif ch == "{":
                depth += 1
            elif ch == "}":
                depth -= 1
                if depth == 0:
                    return text[start:i + 1]
        return ""  # 未配平（截断的 JSON），交由上层兜底

    def _validate_extraction(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """验证提取结果的完整性和合理性"""
        warnings = []
        scores = {}

        # 检查个人信息
        personal = data.get("personal_info", {})
        if not personal.get("name"):
            warnings.append("缺少姓名")
        if not personal.get("highest_education"):
            warnings.append("缺少最高学历")

        # 检查教育经历
        edu = data.get("education_history", [])
        if not edu:
            warnings.append("缺少教育经历")
        else:
            for i, e in enumerate(edu):
                if not isinstance(e, dict):
                    continue
                if not e.get("school"):
                    warnings.append(f"教育经历{i+1}缺少学校名")
                if not e.get("degree"):
                    warnings.append(f"教育经历{i+1}缺少学历等级")

        # 检查技能
        skills = data.get("skills", [])
        if not skills:
            warnings.append("缺少技术栈信息")

        # 计算完整度分数
        total_fields = 0
        filled_fields = 0
        for key, value in personal.items():
            total_fields += 1
            if value is not None and value != "":
                filled_fields += 1

        completeness = filled_fields / max(total_fields, 1)
        scores["personal_completeness"] = round(completeness, 2)
        scores["education_count"] = len(edu)
        scores["skills_count"] = len(skills)
        scores["work_count"] = len(data.get("work_experiences", []))
        scores["project_count"] = len(data.get("projects", []))

        return {
            "valid": len(warnings) == 0 or (personal.get("name") and edu),
            "warnings": warnings,
            "scores": scores
        }

    def _save_to_database(self, data: Dict[str, Any], raw_text: str = "") -> Optional[int]:
        """将提取的结构化数据写入SQL数据库（含原始文本和动态扩展属性）"""
        from backend.database.models import hr_db

        try:
            personal = data.get("personal_info") or {}
            if not isinstance(personal, dict):
                personal = {}

            # 字段归一化：把 LLM 抽取的自由文本 job_status / highest_education
            # 映射为稳定枚举，便于后续精确过滤（在读/非全日制等细节由
            # education_history 承载，不在主字段丢失语义）。
            from backend.utils.field_norm import (
                normalize_job_status,
                normalize_education_level,
            )
            norm_job_status = normalize_job_status(personal.get("job_status"))
            norm_education = normalize_education_level(personal.get("highest_education"))

            # 补全 age：优先从教育经历最早入学时间推算，兜底用 work_years
            age = personal.get("age")
            if not age:
                age = self._estimate_age_from_education(
                    data.get("education_history", []),
                    personal.get("work_years")
                )

            # 1. 写入candidates主表
            candidate_id = hr_db.insert_candidate({
                "name": personal.get("name", "未知"),
                "gender": personal.get("gender"),
                "birth_date": personal.get("birth_date"),
                "age": age,
                "phone": personal.get("phone"),
                "email": personal.get("email"),
                "address": personal.get("address"),
                "current_city": personal.get("current_city"),
                "hometown": personal.get("hometown"),
                "current_position": personal.get("current_position"),
                "work_years": personal.get("work_years", 0),
                "current_salary": personal.get("current_salary"),
                "expected_salary": personal.get("expected_salary"),
                "job_status": norm_job_status,
                "location": personal.get("location") or personal.get("city"),
                "highest_education": norm_education,
                "summary": personal.get("summary"),
            })

            # 2. 写入教育经历表（单条失败不影响其它）
            # 院校层级用权威知识库强制归一化：覆盖 LLM 主观判断，国科大算985，
            # 名单外学校保留 LLM 判断作兜底（详见 backend/utils/school_tier.py）。
            from backend.utils.school_tier import normalize_school_tier
            for edu in data.get("education_history", []) or []:
                if not isinstance(edu, dict):
                    continue
                try:
                    authoritative_tier = normalize_school_tier(
                        edu.get("school"), edu.get("school_tier")
                    )
                    hr_db.insert_education_history(candidate_id, {
                        "degree": edu.get("degree"),
                        "school": edu.get("school"),
                        "major": edu.get("major"),
                        "start_date": edu.get("start_date"),
                        "end_date": edu.get("end_date"),
                        "is_fulltime": edu.get("is_fulltime", True),
                        "school_tier": authoritative_tier,
                    })
                except Exception as e:
                    logger.warning(f"插入教育经历失败，跳过: {e}")

            # 3. 写入技术栈表
            for skill in data.get("skills", []) or []:
                try:
                    if isinstance(skill, str):
                        hr_db.insert_skill(candidate_id, skill, 3, None)
                    elif isinstance(skill, dict) and skill.get("skill_name"):
                        hr_db.insert_skill(
                            candidate_id,
                            skill.get("skill_name", ""),
                            skill.get("proficiency", 3),
                            skill.get("category")
                        )
                except Exception as e:
                    logger.warning(f"插入技能失败，跳过: {e}")

            # 4. 写入获奖证书表
            for award in data.get("awards_certificates", []) or []:
                try:
                    if isinstance(award, str):
                        hr_db.insert_award_certificate(candidate_id, {"type": "award", "name": award})
                    elif isinstance(award, dict):
                        hr_db.insert_award_certificate(candidate_id, {
                            "type": award.get("type", "award"),
                            "name": award.get("name"),
                            "level": award.get("level"),
                            "date": award.get("date"),
                            "role": award.get("role"),
                            "description": award.get("description"),
                        })
                except Exception as e:
                    logger.warning(f"插入获奖证书失败，跳过: {e}")

            # 5. 写入工作经历表
            for exp in data.get("work_experiences", []) or []:
                if not isinstance(exp, dict):
                    continue
                try:
                    hr_db.insert_work_experience(candidate_id, {
                        "company_name": exp.get("company_name"),
                        "position": exp.get("position"),
                        "location": exp.get("location"),
                        "start_date": exp.get("start_date"),
                        "end_date": exp.get("end_date"),
                        "duration_months": exp.get("duration_months"),
                        "description": exp.get("description"),
                    })
                except Exception as e:
                    logger.warning(f"插入工作经历失败，跳过: {e}")

            # 6. 写入项目经历表
            for proj in data.get("projects", []) or []:
                if not isinstance(proj, dict):
                    continue
                try:
                    hr_db.insert_project(candidate_id, {
                        "project_name": proj.get("project_name"),
                        "role": proj.get("role"),
                        "start_date": proj.get("start_date"),
                        "end_date": proj.get("end_date"),
                        "duration_months": proj.get("duration_months"),
                        "description": proj.get("description"),
                        "technologies": proj.get("technologies"),
                    })
                except Exception as e:
                    logger.warning(f"插入项目经历失败，跳过: {e}")

            # 7. 写入论文发表表
            for pub in data.get("publications", []) or []:
                if not isinstance(pub, dict):
                    continue
                try:
                    hr_db.insert_publication(candidate_id, {
                        "title": pub.get("title"),
                        "venue": pub.get("venue"),
                        "venue_type": pub.get("venue_type"),
                        "venue_rank": pub.get("venue_rank"),
                        "sci_zone": pub.get("sci_zone"),
                        "year": pub.get("year"),
                        "authors": pub.get("authors"),
                        "author_position": pub.get("author_position"),
                        "doi": pub.get("doi"),
                        "abstract": pub.get("abstract"),
                    })
                except Exception as e:
                    logger.warning(f"插入论文发表失败，跳过: {e}")

            # 8. 写入国际会议表
            for conf in data.get("conferences", []) or []:
                if not isinstance(conf, dict):
                    continue
                try:
                    hr_db.insert_conference(candidate_id, {
                        "conference_name": conf.get("conference_name"),
                        "conference_rank": conf.get("conference_rank"),
                        "year": conf.get("year"),
                        "location": conf.get("location"),
                        "role": conf.get("role"),
                        "paper_title": conf.get("paper_title"),
                        "description": conf.get("description"),
                    })
                except Exception as e:
                    logger.warning(f"插入国际会议失败，跳过: {e}")

            # 9. 存储简历原始文本
            if raw_text:
                try:
                    hr_db.update_resume_raw_text(candidate_id, raw_text)
                except Exception as e:
                    logger.warning(f"存储简历原文失败，跳过: {e}")

            # 10. 写入动态扩展属性
            extra_attrs = data.get("extra_attributes")
            if extra_attrs and isinstance(extra_attrs, dict):
                try:
                    # 过滤掉 None 值
                    valid_attrs = {k: v for k, v in extra_attrs.items() if v is not None}
                    if valid_attrs:
                        hr_db.insert_extra_attributes_batch(candidate_id, valid_attrs)
                except Exception as e:
                    logger.warning(f"插入动态扩展属性失败，跳过: {e}")

            logger.info(f"Successfully saved candidate {candidate_id} to database")
            return candidate_id

        except Exception as e:
            logger.error(f"Failed to save to database: {e}")
            return None

    @staticmethod
    def _estimate_age_from_education(education_history: list, work_years) -> int:
        """从教育经历推算年龄

        策略：找到最早的入学时间，假设入学时18岁，推算当前年龄。
        兜底：用 22 + work_years。
        """
        from datetime import date

        current_year = date.today().year

        if education_history:
            earliest_start_year = None
            for edu in education_history or []:
                if not isinstance(edu, dict):
                    continue
                start_date = str(edu.get("start_date", ""))
                if not start_date:
                    continue
                match = re.match(r"(\d{4})", start_date)
                if match:
                    year = int(match.group(1))
                    if earliest_start_year is None or year < earliest_start_year:
                        earliest_start_year = year

            if earliest_start_year:
                return current_year - earliest_start_year + 18

        # 兜底
        try:
            return int(22 + float(work_years or 0))
        except (ValueError, TypeError):
            return 22

    async def extract_and_index(self, file_path: str = "", resume_text: str = "") -> Dict[str, Any]:
        """提取简历并同时建立向量索引（完整流程）

        支持两种输入：
            file_path: 简历文件路径
            resume_text: 直接传入的简历纯文本（优先）
        """
        # Step 1: 提取并写入SQL（优先使用直接传入的文本）
        if resume_text and resume_text.strip():
            result = await self.execute({"resume_text": resume_text, "save_to_db": True})
        else:
            result = await self.execute({"file_path": file_path, "save_to_db": True})
        if not result["success"]:
            return result

        candidate_id = result["candidate_id"]
        extracted_data = result["extracted_data"]
        # 获取原始简历文本（用于向量化和存储到向量库 document 字段）
        original_raw_text = result.get("raw_text", "")

        # Step 2: 生成向量并写入ChromaDB
        try:
            from backend.models.multimodal_fusion import multimodal_fusion
            from backend.vector_db.client import vector_db

            # 构建结构化简历文本用于向量化（语义更紧凑）
            structured_text = self._build_resume_text(extracted_data)

            # 生成BGE-M3向量（用结构化文本，语义更集中）
            embedding = multimodal_fusion.extract_text_features(structured_text).flatten().tolist()

            # 构建metadata（元素可能是 str 或 dict，需做类型守卫）
            personal = extracted_data.get("personal_info") or {}
            if not isinstance(personal, dict):
                personal = {}

            def _skill_name(s):
                if isinstance(s, dict):
                    return s.get("skill_name") or s.get("name") or ""
                return str(s) if s else ""

            skills_list = [_skill_name(s) for s in extracted_data.get("skills", [])]
            skills_list = [s for s in skills_list if s]
            schools_list = [
                e.get("school", "")
                for e in extracted_data.get("education_history", [])
                if isinstance(e, dict)
            ]
            schools_list = [s for s in schools_list if s]

            # 动态扩展属性也写入 metadata（便于过滤）
            extra_attrs = extracted_data.get("extra_attributes") or {}
            if not isinstance(extra_attrs, dict):
                extra_attrs = {}

            metadata = {
                "name": personal.get("name", ""),
                "highest_education": personal.get("highest_education", ""),
                "work_years": personal.get("work_years", 0),
                "current_position": personal.get("current_position", ""),
                "location": personal.get("location", ""),
                "skills_text": ",".join(skills_list),
                "school_list": ",".join(schools_list),
            }
            # 将非空的动态属性加入 metadata
            for attr_key, attr_val in extra_attrs.items():
                if attr_val is not None and attr_val != "":
                    metadata[f"extra_{attr_key}"] = str(attr_val) if not isinstance(attr_val, (str, int, float, bool)) else attr_val

            # 写入ChromaDB（document 存储原始简历文本，便于前端展示和 BM25 回溯）
            vector_db.add_candidate(
                candidate_id=candidate_id,
                embedding=embedding,
                metadata=metadata,
                document=original_raw_text or structured_text
            )

            result["vector_indexed"] = True
            logger.info(f"Vector index created for candidate {candidate_id}")

        except Exception as e:
            logger.error(f"Vector indexing failed: {e}")
            result["vector_indexed"] = False
            result["vector_error"] = str(e)

        return result

    def _build_resume_text(self, data: Dict[str, Any]) -> str:
        """从提取的结构化数据构建完整简历文本（用于向量化）"""
        sections = []
        personal = data.get("personal_info") or {}
        if not isinstance(personal, dict):
            personal = {}

        # 个人信息段
        personal_parts = []
        if personal.get("name"):
            personal_parts.append(personal["name"])
        if personal.get("gender"):
            personal_parts.append(personal["gender"])
        if personal.get("age"):
            personal_parts.append(f"{personal['age']}岁")
        if personal.get("current_position"):
            personal_parts.append(personal["current_position"])
        if personal.get("work_years"):
            personal_parts.append(f"{personal['work_years']}年工作经验")
        if personal.get("location"):
            personal_parts.append(f"坐标{personal['location']}")
        if personal.get("job_status"):
            personal_parts.append(personal["job_status"])
        if personal.get("expected_salary"):
            personal_parts.append(f"期望薪资{personal['expected_salary']}元/月")
        sections.append(f"个人简历：{'，'.join(personal_parts)}。")

        # 教育经历段
        edu_list = [e for e in (data.get("education_history") or []) if isinstance(e, dict)]
        if edu_list:
            edu_lines = []
            for edu in edu_list:
                start = edu.get("start_date", "")
                end = edu.get("end_date", "")
                school = edu.get("school", "")
                major = edu.get("major", "")
                degree = edu.get("degree", "")
                tier = edu.get("school_tier", "")
                ft = "全日制" if edu.get("is_fulltime", True) else "非全日制"
                tier_str = f"{tier}/" if tier else ""
                edu_lines.append(f"{start}-{end}，{school}，{major}，{degree}（{tier_str}{ft}）")
            sections.append(f"教育经历：{'；'.join(edu_lines)}。")

        # 技术栈段
        skills = [s for s in (data.get("skills") or []) if isinstance(s, dict) and s.get("skill_name")]
        if skills:
            proficiency_map = {5: "精通", 4: "熟练", 3: "熟悉", 2: "掌握", 1: "了解"}
            sorted_skills = sorted(skills, key=lambda s: s.get("proficiency", 0), reverse=True)
            skill_strs = [f"{s['skill_name']}({proficiency_map.get(s.get('proficiency', 3), '熟悉')})"
                          for s in sorted_skills]
            sections.append(f"技术栈：{'、'.join(skill_strs)}。")

        # 获奖证书段
        awards = [a for a in (data.get("awards_certificates") or []) if isinstance(a, dict)]
        if awards:
            award_lines = [f"{a.get('name', '')}({a.get('level', '')})" for a in awards]
            sections.append(f"获奖证书：{'、'.join(award_lines)}。")

        # 工作经历段
        work_exps = [w for w in (data.get("work_experiences") or []) if isinstance(w, dict)]
        if work_exps:
            exp_lines = []
            for exp in work_exps:
                company = exp.get("company_name", "")
                pos = exp.get("position", "")
                duration = exp.get("duration_months")
                desc = exp.get("description", "")
                duration_str = f"({duration}个月)" if duration else ""
                desc_str = f"-{desc}" if desc else ""
                exp_lines.append(f"{company}-{pos}{duration_str}{desc_str}")
            sections.append(f"工作经历：{'；'.join(exp_lines)}。")

        # 项目经历段
        projects = [p for p in data.get("projects", []) if isinstance(p, dict)]
        if projects:
            proj_lines = []
            for proj in projects:
                name = proj.get("project_name", "")
                role = proj.get("role", "")
                techs = proj.get("technologies", "")
                tech_str = f"({techs})" if techs else ""
                proj_lines.append(f"{name}{tech_str}-{role}")
            sections.append(f"项目经历：{'；'.join(proj_lines)}。")

        # 论文发表段
        pubs = [p for p in (data.get("publications") or []) if isinstance(p, dict)]
        if pubs:
            pub_lines = []
            pos_map = {"first_author": "第一作者", "corresponding_author": "通讯作者", "co_author": "合作者"}
            for pub in pubs:
                title = pub.get("title", "")
                venue = pub.get("venue", "")
                rank = pub.get("venue_rank", "")
                sci = f"/SCI {pub['sci_zone']}" if pub.get("sci_zone") else ""
                year = pub.get("year", "")
                pos = pos_map.get(pub.get("author_position", ""), "")
                pub_lines.append(f"{title}({venue},{rank}{sci},{year},{pos})")
            sections.append(f"论文发表：{'；'.join(pub_lines)}。")

        # 国际会议段
        confs = [c for c in (data.get("conferences") or []) if isinstance(c, dict)]
        if confs:
            conf_lines = []
            role_map = {
                "oral_presentation": "口头报告",
                "poster": "海报展示",
                "workshop": "研讨会",
                "attendee": "参会者",
                "session_chair": "分会主席",
                "keynote": "特邀报告",
            }
            for conf in confs:
                name = conf.get("conference_name", "")
                rank = conf.get("conference_rank", "")
                year = conf.get("year", "")
                role = role_map.get(conf.get("role", ""), conf.get("role", ""))
                paper = f",论文:{conf['paper_title']}" if conf.get("paper_title") else ""
                conf_lines.append(f"{name}({rank},{year},{role}{paper})")
            sections.append(f"国际会议：{'；'.join(conf_lines)}。")

        return "\n".join(sections)
